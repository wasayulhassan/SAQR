"""
ai_ppt.py
Turns a file attached in the Chat panel + a short guided Q&A into a
custom-designed PowerPoint deck: the model proposes its own outline and
explains its choices, real charts get generated from the file's actual
data where relevant, a few slides can get AI-generated illustrations, and
the whole thing is styled with a randomly/answer-picked visual theme so
decks don't all look the same.

This is deliberately separate from ppt_gen.py (the Export tab's fixed
report-style deck) and chart_builder.py (manual one-off chart requests) —
this one is the "have a conversation about my file, then design something
around it" path.
"""

import os
import json
import re
import uuid
import random

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

from . import chatbot, analyzer, chart_builder, image_gen, ppt_themes

OUTPUT_DIR = os.path.join(os.path.dirname(__file__), "..", "outputs")
GEN_IMG_DIR = os.path.join(os.path.dirname(__file__), "..", "static", "generated")
os.makedirs(OUTPUT_DIR, exist_ok=True)
os.makedirs(GEN_IMG_DIR, exist_ok=True)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


# ---------------------------------------------------------------- outline --

def _outline_system_prompt(file_ctx: dict, answers: dict) -> str:
    is_spreadsheet = file_ctx.get("type") == "spreadsheet"
    slide_count_answer = (answers.get("slide_count") or "auto").strip().lower()
    if slide_count_answer in ("auto", "let saqr decide", ""):
        slide_count_instruction = "Choose an appropriate slide count yourself — typically 6 to 10."
    else:
        slide_count_instruction = f"Aim for {slide_count_answer} content slides (not counting title/closing)."

    columns_note = ""
    if is_spreadsheet:
        columns_note = (
            "\nThe file is a spreadsheet. Only request a chart slide if a chart would "
            "genuinely help — and only use column names that actually exist in the data "
            "(they're listed below in the file content)."
        )

    return f"""You are Saqr, an assistant that designs custom presentations from a
student's or professional's uploaded file. Given the file's content and the
user's answers about purpose, tone, and style, propose your OWN outline —
don't just restate the file, add structure, insight, and opinion like a
good editor would.

{slide_count_instruction}{columns_note}

Respond with ONLY valid JSON, no markdown code fences, matching exactly this
shape:
{{
  "title": "deck title",
  "subtitle": "short subtitle",
  "rationale": "2-4 sentences, first person, explaining how you structured this and why — your own opinion on what mattered most",
  "slides": [
    {{"type": "content", "heading": "...", "bullets": ["...", "..."]}},
    {{"type": "chart", "heading": "...", "bullets": ["one short line of context"], "chart_type": "bar|column|line|pie|scatter|horizontal_bar", "x_col": "<exact column name>", "y_cols": ["<exact column name>"]}},
    {{"type": "image", "heading": "...", "bullets": ["one short caption"], "image_prompt": "a short, safe, descriptive prompt for a text-to-image model, no real people, no copyrighted characters"}}
  ],
  "closing": "closing slide heading, e.g. a thank-you or call to action"
}}

Only include "type": "chart" slides if the file is a spreadsheet with
relevant numeric columns. Use "type": "image" sparingly (at most 2-3 across
the whole deck) for slides that would benefit from an illustrative visual.
Everything else should be "type": "content"."""


def _outline_user_prompt(file_ctx: dict, answers: dict) -> str:
    parts = [
        f"Purpose / audience: {answers.get('purpose') or 'not specified'}",
        f"Tone: {answers.get('tone') or 'not specified'}",
        f"Visual style requested: {answers.get('style') or 'not specified'}",
    ]
    if answers.get("extra"):
        parts.append(f"Extra instructions: {answers['extra']}")
    parts.append(
        f"\nAttached file: {file_ctx.get('filename')} ({file_ctx.get('meta')})\n"
        f"Extracted content:\n{file_ctx.get('content_text', '')}"
    )
    return "\n".join(parts)


def _strip_code_fence(text: str) -> str:
    text = text.strip()
    text = re.sub(r"^```(?:json)?\s*", "", text)
    text = re.sub(r"\s*```$", "", text)
    return text.strip()


def _fallback_outline(file_ctx: dict, answers: dict) -> dict:
    """Used if the hosted model is unavailable or returns something we
    can't parse — the feature should still produce a usable deck."""
    filename = file_ctx.get("filename", "your file")
    slides = [
        {"type": "content", "heading": "About this file",
         "bullets": [f"Source: {filename}", f"{file_ctx.get('meta', '')}"]},
        {"type": "content", "heading": "Key content",
         "bullets": [line.strip() for line in file_ctx.get("content_text", "").splitlines() if line.strip()][:6]
         or ["Content extracted from the attached file."]},
    ]
    return {
        "title": answers.get("purpose") or f"Presentation: {filename}",
        "subtitle": "Generated by SAQR",
        "rationale": (
            "I wasn't able to reach the hosted model for a fully custom outline, so "
            "this is a straightforward structure built directly from your file's "
            "content — still worth a look, but feel free to ask me to try again."
        ),
        "slides": slides,
        "closing": "Thank you",
    }


def build_outline(file_ctx: dict, answers: dict) -> dict:
    try:
        raw = chatbot.chat_raw(
            _outline_system_prompt(file_ctx, answers),
            _outline_user_prompt(file_ctx, answers),
            max_tokens=1600,
        )
        outline = json.loads(_strip_code_fence(raw))
        if not isinstance(outline, dict) or "slides" not in outline:
            raise ValueError("missing 'slides'")
        return outline
    except Exception:
        return _fallback_outline(file_ctx, answers)


# ------------------------------------------------------------- pptx build --

def _new_slide(prs, theme):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = theme["bg"]
    return slide


def _textbox(slide, left, top, width, height):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    return box, tf


def _accent_bar(slide, theme, left=Inches(0), top=Inches(0), width=Inches(0.18), height=SLIDE_H):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    bar.fill.solid()
    bar.fill.fore_color.rgb = theme["accent"]
    bar.line.fill.background()
    return bar


def _add_title_slide(prs, theme, title, subtitle):
    slide = _new_slide(prs, theme)
    _accent_bar(slide, theme)

    box, tf = _textbox(slide, Inches(0.9), Inches(2.6), Inches(11.5), Inches(1.8))
    p = tf.paragraphs[0]
    p.text = title or "Presentation"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.name = theme["heading_font"]
    p.font.color.rgb = theme["text"]

    box2, tf2 = _textbox(slide, Inches(0.95), Inches(3.9), Inches(11), Inches(1))
    p2 = tf2.paragraphs[0]
    p2.text = subtitle or ""
    p2.font.size = Pt(18)
    p2.font.name = theme["body_font"]
    p2.font.color.rgb = theme["muted"]
    return slide


def _add_content_slide(prs, theme, heading, bullets):
    slide = _new_slide(prs, theme)
    _accent_bar(slide, theme)

    _, tf = _textbox(slide, Inches(0.9), Inches(0.7), Inches(11.5), Inches(1))
    p = tf.paragraphs[0]
    p.text = heading or ""
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.name = theme["heading_font"]
    p.font.color.rgb = theme["text"]

    _, body_tf = _textbox(slide, Inches(1.1), Inches(1.9), Inches(10.8), Inches(5))
    for i, b in enumerate(bullets or []):
        bp = body_tf.paragraphs[0] if i == 0 else body_tf.add_paragraph()
        bp.text = f"•  {b}"
        bp.font.size = Pt(18)
        bp.font.name = theme["body_font"]
        bp.font.color.rgb = theme["text"]
        bp.space_after = Pt(14)
    return slide


def _add_chart_slide(prs, theme, heading, bullets, chart_info, df):
    chart_path = None
    if df is not None and chart_info:
        result = chart_builder.build_chart(
            df,
            chart_type=chart_info.get("chart_type", "bar"),
            x_col=chart_info.get("x_col", ""),
            y_cols=chart_info.get("y_cols", []),
            title=heading,
            color="#%02X%02X%02X" % (theme["accent"][0], theme["accent"][1], theme["accent"][2]),
        )
        if result.get("ok"):
            chart_path = result["path"]

    if not chart_path:
        return _add_content_slide(prs, theme, heading, bullets or ["(chart could not be generated)"])

    slide = _new_slide(prs, theme)
    _accent_bar(slide, theme)

    _, tf = _textbox(slide, Inches(0.9), Inches(0.6), Inches(11.5), Inches(0.9))
    p = tf.paragraphs[0]
    p.text = heading or ""
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.name = theme["heading_font"]
    p.font.color.rgb = theme["text"]

    slide.shapes.add_picture(chart_path, Inches(2.4), Inches(1.55), height=Inches(4.9))

    if bullets:
        _, cap_tf = _textbox(slide, Inches(1.1), Inches(6.9), Inches(11), Inches(0.5))
        cp = cap_tf.paragraphs[0]
        cp.text = bullets[0]
        cp.font.size = Pt(12)
        cp.font.italic = True
        cp.font.color.rgb = theme["muted"]
    return slide


def _add_image_slide(prs, theme, heading, bullets, image_prompt, images_used):
    img_path = None
    if image_prompt and images_used[0] < image_gen.MAX_IMAGES_PER_DECK and image_gen.is_configured():
        candidate_path = os.path.join(GEN_IMG_DIR, f"gen_{uuid.uuid4().hex[:10]}.png")
        if image_gen.generate_image(image_prompt, candidate_path):
            img_path = candidate_path
            images_used[0] += 1

    slide = _new_slide(prs, theme)
    _accent_bar(slide, theme)

    _, tf = _textbox(slide, Inches(0.9), Inches(0.6), Inches(11.5), Inches(0.9))
    p = tf.paragraphs[0]
    p.text = heading or ""
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.name = theme["heading_font"]
    p.font.color.rgb = theme["text"]

    if img_path:
        slide.shapes.add_picture(img_path, Inches(3.4), Inches(1.6), height=Inches(5))
    else:
        # graceful fallback: a themed decorative block instead of a broken slide
        block = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(3.4), Inches(1.8), Inches(6.5), Inches(4.4))
        block.fill.solid()
        block.fill.fore_color.rgb = theme["panel"]
        block.line.color.rgb = theme["accent"]
        block.line.width = Pt(1.5)
        block.text_frame.word_wrap = True
        bp = block.text_frame.paragraphs[0]
        bp.text = (bullets[0] if bullets else heading) or ""
        bp.font.size = Pt(16)
        bp.font.color.rgb = theme["muted"]
        bp.alignment = PP_ALIGN.CENTER
        block.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

    if bullets:
        _, cap_tf = _textbox(slide, Inches(1.1), Inches(6.9), Inches(11), Inches(0.5))
        cp = cap_tf.paragraphs[0]
        cp.text = bullets[0]
        cp.font.size = Pt(12)
        cp.font.italic = True
        cp.font.color.rgb = theme["muted"]
    return slide


def _add_closing_slide(prs, theme, text):
    slide = _new_slide(prs, theme)
    band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_W, SLIDE_H)
    band.fill.solid()
    band.fill.fore_color.rgb = theme["accent"]
    band.line.fill.background()

    _, tf = _textbox(slide, Inches(1), Inches(3.2), Inches(11.3), Inches(1.2))
    p = tf.paragraphs[0]
    p.text = text or "Thank you"
    p.font.size = Pt(34)
    p.font.bold = True
    p.font.name = theme["heading_font"]
    p.font.color.rgb = theme["bg"] if theme["key"] != "brand" else theme["text"]
    p.alignment = PP_ALIGN.CENTER
    return slide


def generate_presentation(outline: dict, theme: dict, file_ctx: dict, saved_filepath: str = None) -> str:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    _add_title_slide(prs, theme, outline.get("title"), outline.get("subtitle"))

    df = None
    if file_ctx.get("type") == "spreadsheet" and saved_filepath and os.path.exists(saved_filepath):
        try:
            df = analyzer.load_file(saved_filepath)
        except Exception:
            df = None

    images_used = [0]  # mutable counter passed by reference

    for slide_def in outline.get("slides", []):
        stype = slide_def.get("type", "content")
        heading = slide_def.get("heading", "")
        bullets = slide_def.get("bullets", [])

        if stype == "chart":
            chart_info = {
                "chart_type": slide_def.get("chart_type", "bar"),
                "x_col": slide_def.get("x_col", ""),
                "y_cols": slide_def.get("y_cols", []),
            }
            _add_chart_slide(prs, theme, heading, bullets, chart_info, df)
        elif stype == "image":
            _add_image_slide(prs, theme, heading, bullets, slide_def.get("image_prompt", ""), images_used)
        else:
            _add_content_slide(prs, theme, heading, bullets)

    _add_closing_slide(prs, theme, outline.get("closing", "Thank you"))

    fname = f"ai_deck_{uuid.uuid4().hex[:8]}.pptx"
    fpath = os.path.join(OUTPUT_DIR, fname)
    prs.save(fpath)
    return fpath
