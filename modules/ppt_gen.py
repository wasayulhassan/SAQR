"""
ppt_gen.py
Builds a PowerPoint (.pptx) deck from an analysis dict produced by analyzer.py
"""

import os
import uuid
from datetime import datetime
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

OUTPUT_DIR = os.path.join(os.path.dirname(__file__), "..", "outputs")
os.makedirs(OUTPUT_DIR, exist_ok=True)

ACCENT = RGBColor(0x4F, 0x46, 0xE5)


def _add_title_slide(prs, title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle
    return slide


def _add_bullet_slide(prs, title, bullets):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    body = slide.placeholders[1].text_frame
    body.clear()
    for i, b in enumerate(bullets):
        p = body.paragraphs[0] if i == 0 else body.add_paragraph()
        p.text = b
        p.font.size = Pt(18)
    return slide


def _add_image_slide(prs, title, image_path):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = title
    if os.path.exists(image_path):
        slide.shapes.add_picture(image_path, Inches(1.2), Inches(1.3), width=Inches(7.5))
    return slide


def build_presentation(analysis: dict, title: str = "Data Analysis") -> str:
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    summary = analysis["summary"]
    trends = analysis["trends"]
    anomalies = analysis["anomalies"]
    charts = analysis["charts"]

    _add_title_slide(prs, title, f"Generated {datetime.now().strftime('%Y-%m-%d')}")

    _add_bullet_slide(prs, "Overview", [
        f"Rows: {summary['rows']}",
        f"Columns: {len(summary['columns'])}",
        f"Numeric columns: {len(summary['numeric_columns'])}",
    ])

    trend_bullets = []
    for col, t in trends.items():
        if isinstance(t, dict):
            trend_bullets.append(f"{col}: {t['direction']} ({t['pct_change_start_to_end']}%)")
        else:
            trend_bullets.append(f"{col}: {t}")
    if trend_bullets:
        _add_bullet_slide(prs, "Trends", trend_bullets[:8])

    if anomalies:
        anomaly_bullets = [
            f"{col}: {len(info['outlier_row_indices'])} outlier(s)"
            for col, info in anomalies.items()
        ]
        _add_bullet_slide(prs, "Anomalies", anomaly_bullets)
    else:
        _add_bullet_slide(prs, "Anomalies", ["No significant anomalies detected."])

    for chart_path in charts:
        chart_name = os.path.basename(chart_path)
        _add_image_slide(prs, "Chart", chart_path)

    fname = f"deck_{uuid.uuid4().hex[:8]}.pptx"
    fpath = os.path.join(OUTPUT_DIR, fname)
    prs.save(fpath)
    return fpath
